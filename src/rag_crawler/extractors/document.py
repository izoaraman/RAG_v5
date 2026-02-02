"""
Document extractors for PDF, DOCX, PPTX, XLSX files.

Downloads and extracts text content from binary document formats.
"""

import io
import tempfile
from pathlib import Path
from urllib.parse import urlparse

import httpx
from pypdf import PdfReader


class DocumentExtractor:
    """
    Extracts text from document files (PDF, DOCX, PPTX, XLSX).

    Downloads the document and extracts text using appropriate libraries.
    """

    def __init__(self, timeout: float = 30.0):
        self.timeout = timeout
        self.client = httpx.AsyncClient(timeout=timeout, follow_redirects=True)

    async def close(self) -> None:
        """Close the HTTP client."""
        await self.client.aclose()

    async def __aenter__(self):
        return self

    async def __aexit__(self, exc_type, exc_val, exc_tb):
        await self.close()

    async def extract(self, url: str) -> tuple[str, str]:
        """
        Download and extract text from a document URL.

        Args:
            url: URL of the document to extract.

        Returns:
            Tuple of (extracted_text, document_type).

        Raises:
            ValueError: If document type is not supported.
            httpx.HTTPError: If download fails.
        """
        doc_type = self._get_document_type(url)
        if not doc_type:
            raise ValueError(f"Unsupported document type for URL: {url}")

        content = await self._download(url)

        if doc_type == "pdf":
            text = self._extract_pdf(content)
        elif doc_type in ("docx", "doc"):
            text = self._extract_docx(content)
        elif doc_type in ("pptx", "ppt"):
            text = self._extract_pptx(content)
        elif doc_type in ("xlsx", "xls"):
            text = self._extract_xlsx(content)
        elif doc_type in ("txt", "csv", "rtf"):
            text = self._extract_text(content)
        else:
            text = f"[Document type '{doc_type}' extraction not implemented]"

        return text, doc_type

    def _get_document_type(self, url: str) -> str | None:
        """Determine document type from URL."""
        parsed = urlparse(url)
        path = parsed.path.lower()

        extensions = {
            ".pdf": "pdf",
            ".docx": "docx",
            ".doc": "doc",
            ".pptx": "pptx",
            ".ppt": "ppt",
            ".xlsx": "xlsx",
            ".xls": "xls",
            ".txt": "txt",
            ".csv": "csv",
            ".rtf": "rtf",
            ".odt": "odt",
            ".odp": "odp",
            ".ods": "ods",
        }

        for ext, doc_type in extensions.items():
            if path.endswith(ext):
                return doc_type
        return None

    async def _download(self, url: str) -> bytes:
        """Download document content."""
        response = await self.client.get(url)
        response.raise_for_status()
        return response.content

    def _extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF content."""
        try:
            reader = PdfReader(io.BytesIO(content))
            text_parts = []

            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"## Page {i + 1}\n\n{page_text}")

            return "\n\n".join(text_parts) if text_parts else "[No text extracted from PDF]"
        except Exception as e:
            return f"[PDF extraction error: {e}]"

    def _extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX content."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            paragraphs = []

            for para in doc.paragraphs:
                if para.text.strip():
                    # Detect headings
                    if para.style and para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        try:
                            heading_level = int(level)
                            prefix = "#" * min(heading_level, 6)
                            paragraphs.append(f"{prefix} {para.text}")
                        except ValueError:
                            paragraphs.append(para.text)
                    else:
                        paragraphs.append(para.text)

            # Extract tables
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append("| " + " | ".join(cells) + " |")
                if table_rows:
                    paragraphs.append("\n".join(table_rows))

            return "\n\n".join(paragraphs) if paragraphs else "[No text extracted from DOCX]"
        except Exception as e:
            return f"[DOCX extraction error: {e}]"

    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PPTX content."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            slides_text = []

            for i, slide in enumerate(prs.slides):
                slide_parts = [f"## Slide {i + 1}"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)

                slides_text.append("\n\n".join(slide_parts))

            return "\n\n---\n\n".join(slides_text) if slides_text else "[No text extracted from PPTX]"
        except Exception as e:
            return f"[PPTX extraction error: {e}]"

    def _extract_xlsx(self, content: bytes) -> str:
        """Extract text from XLSX content."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), data_only=True)
            sheets_text = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_parts = [f"## Sheet: {sheet_name}"]

                rows = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    if any(cells):  # Skip empty rows
                        rows.append("| " + " | ".join(cells) + " |")

                if rows:
                    sheet_parts.append("\n".join(rows))
                    sheets_text.append("\n\n".join(sheet_parts))

            return "\n\n---\n\n".join(sheets_text) if sheets_text else "[No text extracted from XLSX]"
        except Exception as e:
            return f"[XLSX extraction error: {e}]"

    def _extract_text(self, content: bytes) -> str:
        """Extract text from plain text files."""
        try:
            # Try UTF-8 first, then fallback to latin-1
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                return content.decode("latin-1")
        except Exception as e:
            return f"[Text extraction error: {e}]"
